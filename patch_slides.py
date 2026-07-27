#!/usr/bin/env python3
"""
Targeted patch for HPC - xCAT Overview - Roary.pptx
Fixes wrong IPs/MACs, makedns syntax, and slide 24 content.
Does NOT touch slides 15, 16, 19 (user manually improved those).
"""
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

FILE = "HPC - xCAT Overview - Roary.pptx"


# ── Helpers (same as update_slides.py) ──────────────────────────────────────

def _replace_in_para(para, old, new):
    for run in para.runs:
        if old in run.text:
            run.text = run.text.replace(old, new)
            return
    full = "".join(r.text for r in para.runs)
    if old in full:
        if para.runs:
            para.runs[0].text = full.replace(old, new)
            for r in para.runs[1:]:
                r.text = ""

def _replace_in_tf(tf, old, new):
    for para in tf.paragraphs:
        _replace_in_para(para, old, new)

def _replace_in_shape(shape, old, new):
    if shape.has_text_frame:
        _replace_in_tf(shape.text_frame, old, new)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_in_tf(cell.text_frame, old, new)
    if shape.shape_type == 6:  # GROUP
        for child in shape.shapes:
            _replace_in_shape(child, old, new)

def replace_in_slide(slide, old, new):
    for shape in slide.shapes:
        _replace_in_shape(shape, old, new)

def set_tf_text(tf, lines):
    """Overwrite text frame with new lines, preserving first-run formatting."""
    txBody = tf._txBody
    existing = txBody.findall(qn("a:p"))
    tmpl_pPr = tmpl_rPr = None
    if existing:
        pPr = existing[0].find(qn("a:pPr"))
        if pPr is not None:
            tmpl_pPr = copy.deepcopy(pPr)
        first_r = existing[0].find(qn("a:r"))
        if first_r is not None:
            rPr = first_r.find(qn("a:rPr"))
            if rPr is not None:
                tmpl_rPr = copy.deepcopy(rPr)
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    for line in lines:
        p = etree.SubElement(txBody, qn("a:p"))
        if tmpl_pPr is not None:
            p.append(copy.deepcopy(tmpl_pPr))
        r = etree.SubElement(p, qn("a:r"))
        if tmpl_rPr is not None:
            r.append(copy.deepcopy(tmpl_rPr))
        t = etree.SubElement(r, qn("a:t"))
        t.text = line


# ── Slide 20 fixes ───────────────────────────────────────────────────────────

def fix_slide20(slide):
    """Fix wrong gpu-a100-01 IP (246.37→246.171) and MAC in DHCP/DNS examples."""
    pairs = [
        # DHCP query output
        ("192.168.246.37, hardware-address = d4:04:e6:06:04:ac",
         "192.168.246.171, hardware-address = f4:c7:aa:41:65:44"),
        # /etc/hosts entries
        ("192.168.246.37 gpu-a100-01 gpu-a100-01.roary.net",
         "192.168.246.171 gpu-a100-01 gpu-a100-01.roary.net"),
        ("192.168.233.37 gpu-a100-01 gpu-a100-01.roary.net",
         "192.168.233.171 gpu-a100-01 gpu-a100-01.roary.net"),
        # standalone IP in other contexts on this slide
        ("= 192.168.246.37",  "= 192.168.246.171"),
    ]
    for old, new in pairs:
        replace_in_slide(slide, old, new)
    print("  [slide 20] IP/MAC corrections applied")


# ── Slide 24 rewrite ─────────────────────────────────────────────────────────

SLIDE24_LINES = [
    "Adding new compute node/group definitions in xCAT",
    "",
    "Add node to groups:",
    "  nodeadd n126 groups=hpe,alma9              # HPE compute",
    "  nodeadd gpu-a100-01 groups=raptor,alma9    # GPU/Raptor",
    "  nodeadd r126 groups=mgt                    # BMC/DRAC",
    "",
    "HPE compute node (n126) — key attributes:",
    "  chdef -t node n126 ip=192.168.246.196 mac=d4:04:e6:06:04:ac",
    "  chdef -t node n126 netboot=xnba installnic=mac",
    "  chdef -t node n126 os=alma9.5 arch=x86_64 mgt=ipmi bmc=r126",
    "  chdef -t node n126 xcatmaster=hpcxcat",
    '  chdef -t node n126 addkcmdline="biosdevname=0 net.ifnames=0"',
    '  chdef -t node n126 postscripts="syslog,remoteshell,syncfiles,setupntp,',
    "    roary-sethostname,confignetwork,roary-repo-alma9,roary-ldap,",
    '    roary-nouserlogin,roary-setuplustre,roary-slurm"',
    "  chdef -t node n126 postbootscripts=otherpkgs",
    "  (Full bond0/VLAN NIC config: see xCAT Admin Guide)",
    "",
    "GPU node (gpu-a100-01) — key attributes:",
    "  chdef -t node gpu-a100-01 ip=192.168.246.171 mac=f4:c7:aa:41:65:44",
    "  chdef -t node gpu-a100-01 netboot=xnba os=alma9.5 arch=x86_64",
    "  chdef -t node gpu-a100-01 mgt=ipmi bmc=raptor01b xcatmaster=hpcxcat",
    '  chdef -t node gpu-a100-01 postscripts="syslog,remoteshell,syncfiles,',
    "    setupntp,confignetwork,roary-ldap,roary-nouserlogin,",
    '    roary-setuplustre,roary-slurm"',
    "  chdef -t node gpu-a100-01 postbootscripts=otherpkgs",
]

def fix_slide24(slide):
    """Rewrite slide 24 body with accurate, concise node-add commands."""
    # Find the largest text frame (the main body box, not the title)
    largest = None
    largest_len = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text
            if len(t) > largest_len and "nodeadd" in t:
                largest_len = len(t)
                largest = shape
    if largest:
        set_tf_text(largest.text_frame, SLIDE24_LINES)
        print("  [slide 24] body rewritten with accurate commands")
    else:
        print("  [slide 24] WARNING: target shape not found")


# ── Slide 25 fixes ───────────────────────────────────────────────────────────

def fix_slide25(slide):
    """Remove -a flag from makedns and fix n125 → n126."""
    pairs = [
        # makedns does NOT use -a flag
        ("makedns –a gpu-a100-01", "makedns gpu-a100-01"),  # em-dash
        ("makedns –a r-a100-01",  "makedns r-a100-01"),
        ("makedns -a gpu-a100-01",     "makedns gpu-a100-01"),
        ("makedns -a r-a100-01",       "makedns r-a100-01"),
        # Stale n125 reference
        ("makedhcp –a n125",      "makedhcp -a n126"),
        ("makedhcp -a n125",           "makedhcp -a n126"),
        # Remove duplicate trailing makedhcp for gpu-a100-01 if present as standalone
        ("makedhcp –a gpu-a100-01 | makedhcp –a r-a100-01 | makedhcp –a gpu-a100-01",
         "makedhcp -a gpu-a100-01 | makedhcp -a r-a100-01"),
    ]
    for old, new in pairs:
        replace_in_slide(slide, old, new)
    print("  [slide 25] makedns and node name corrections applied")


# ── Slide 26 fixes ───────────────────────────────────────────────────────────

def fix_slide26(slide):
    """Fix stale pxe boot path reference."""
    pairs = [
        ("nodeset for pxe makes changes to /tftpboot/pxelinux.cfg/{node hex ip}",
         "nodeset writes boot config to /tftpboot/xcat/xnba/nodes/"),
        ("Note: nodeset for pxe makes changes to /tftpboot/pxelinux.cfg/{node hex ip}",
         "Note: nodeset writes boot config to /tftpboot/xcat/xnba/nodes/"),
    ]
    for old, new in pairs:
        replace_in_slide(slide, old, new)
    print("  [slide 26] boot path reference updated")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    print(f"Loading {FILE!r} ...")
    prs = Presentation(FILE)
    print(f"  {len(prs.slides)} slides")

    print("Patching slide 20 ...")
    fix_slide20(prs.slides[19])   # 0-indexed

    print("Patching slide 24 ...")
    fix_slide24(prs.slides[23])

    print("Patching slide 25 ...")
    fix_slide25(prs.slides[24])

    print("Patching slide 26 ...")
    fix_slide26(prs.slides[25])

    print(f"Saving {FILE!r} ...")
    prs.save(FILE)
    print("Done.")


if __name__ == "__main__":
    main()
