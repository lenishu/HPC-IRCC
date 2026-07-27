#!/usr/bin/env python3
"""
Update HPC xCAT Overview slides: Panther/CentOS7 -> Roary/AlmaLinux9
Usage:  python update_slides.py
Output: HPC - xCAT Overview - Roary.pptx
"""
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

INPUT_FILE  = "HPC - xCAT Overview.pptx"
OUTPUT_FILE = "HPC - xCAT Overview - Roary.pptx"


# ── Low-level helpers ────────────────────────────────────────────────────────

def _replace_in_para(para, old: str, new: str):
    """Replace text in one paragraph, handling text split across runs."""
    for run in para.runs:
        if old in run.text:
            run.text = run.text.replace(old, new)
            return
    full = "".join(r.text for r in para.runs)
    if old in full:
        if para.runs:
            para.runs[0].text = full.replace(old, new)
            for r in para.runs[1:]:
                r.text = ""


def _replace_in_tf(tf, old: str, new: str):
    for para in tf.paragraphs:
        _replace_in_para(para, old, new)


def _replace_in_shape(shape, old: str, new: str):
    if shape.has_text_frame:
        _replace_in_tf(shape.text_frame, old, new)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                _replace_in_tf(cell.text_frame, old, new)
    # Recurse into grouped shapes
    if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            _replace_in_shape(child, old, new)


def apply_replacements(prs, replacements):
    for old, new in replacements:
        for slide in prs.slides:
            for shape in slide.shapes:
                _replace_in_shape(shape, old, new)


def set_tf_text(tf, lines):
    """Overwrite a text frame with new lines, preserving first-run formatting."""
    txBody = tf._txBody
    existing = txBody.findall(qn("a:p"))
    tmpl_pPr = tmpl_rPr = None
    if existing:
        pPr = existing[0].find(qn("a:pPr"))
        if pPr is not None:
            tmpl_pPr = copy.deepcopy(pPr)
        first_r = existing[0].find(qn("a:r"))
        if first_r is not None:
            rPr = first_r.find(qn("a:rPr"))
            if rPr is not None:
                tmpl_rPr = copy.deepcopy(rPr)
    for p in list(txBody.findall(qn("a:p"))):
        txBody.remove(p)
    for line in lines:
        p = etree.SubElement(txBody, qn("a:p"))
        if tmpl_pPr is not None:
            p.append(copy.deepcopy(tmpl_pPr))
        r = etree.SubElement(p, qn("a:r"))
        if tmpl_rPr is not None:
            r.append(copy.deepcopy(tmpl_rPr))
        t = etree.SubElement(r, qn("a:t"))
        t.text = line


# ── Ordered replacements (specific → broad) ──────────────────────────────────
# Rules: put longer / more specific strings before shorter ones that could
# match the same text.  This avoids double-substitutions.

REPLACEMENTS = [
    # Shell prompt
    ("[[hpcslurm01]root@ms1 ~]#",    "[root@hpcxcat ~]#"),
    ("[[hpcslurm01]root@ms1 dhcp]#", "[root@hpcxcat dhcp]#"),
    ("[hpcslurm01]root@ms1 ~]#",     "[root@hpcxcat ~]#"),
    ("[hpcslurm01]root@ms1 dhcp]#",  "[root@hpcxcat dhcp]#"),

    # Full OS image names  (before any component substitution)
    ("centos7.8-x86_64-netboot-panther-pec6220", "alma9.5-x86_64-netboot-compute"),
    ("centos7.6-x86_64-netboot-panther-pec6220", "alma9.7-x86_64-netboot-compute"),
    ("centos7.8-x86_64-netboot-compute",         "alma9.5-x86_64-netboot-compute"),
    ("centos7.8-x86_64-install-compute",         "alma9.5-x86_64-install-compute"),

    # Full path replacements  (before generic centos→alma)
    ("/install/custom/netboot/centos/panther-pec6220.centos7.x86_64.exlist",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.exlist"),
    ("/install/custom/netboot/centos/panther-pec6220.centos78.x86_64.exlist",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.exlist"),
    ("/install/custom/netboot/centos/panther-pec6220.centos7.x86_64.pkglist",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.pkglist"),
    ("/install/custom/netboot/centos/panther-pec6220.centos78.x86_64.pkglist",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.pkglist"),
    ("/install/custom/netboot/centos/panther-pec6220.centos78.x86_64.otherpkgs.pkglist",
     "/install/custom/netboot/alma/compute.alma9.x86_64.otherpkgs.pkglist"),
    ("/install/custom/netboot/centos/panther-pec6220.centos7.x86_64.postinstall",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.postinstall"),
    ("/install/custom/netboot/centos/panther-pec6220.centos78.x86_64.postinstall",
     "/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.postinstall"),
    ("/install/netboot/centos7.8/x86_64/panther-pec6220",
     "/install/netboot/alma9.5/x86_64/compute"),
    ("/install/netboot/centos7.6/x86_64/compute", "/install/netboot/alma9.7/x86_64/compute"),
    ("/install/netboot/centos7.8/x86_64/compute", "/install/netboot/alma9.5/x86_64/compute"),
    ("/install/centos7.8/x86_64",                 "/install/alma9.5/x86_64"),
    ("/install/centos7.6/x86_64",                 "/install/alma9.7/x86_64"),
    ("/install/post/otherpkgs/centos7.8/x86_64",  "/install/post/otherpkgs/alma9.5/x86_64"),
    ("/install/post/otherpkgs/centos7.6/x86_64",  "/install/post/otherpkgs/alma9.7/x86_64"),

    # mkdef option string
    ("-o centos7.8-x86_64-netboot-panther-pec6220 profile=panther-pec6220",
     "-o alma9.5-x86_64-netboot-compute profile=compute"),

    # Profile
    ("profile=panther-pec6220", "profile=compute"),
    ("panther-pec6220",          "compute"),

    # Postscripts  (specific variants before catch-all)
    ("panther-othernics-centos7",     "roary-othernics-alma9"),
    ("panther-resolv-centos7",        "roary-resolv-alma9"),
    ("panther-yumrepo-centos7",       "roary-yumrepo-alma9"),
    ("panther-rpms-centos7",          "roary-rpms-alma9"),
    ("panther-slurm-centos7_current", "roary-slurm-alma9-current"),
    ("panther-gpucuda11",             "roary-gpucuda"),
    ("panther-",                      "roary-"),       # catch-all

    # Node group combos  (before individual name replacements)
    ("compute-centos78-nomlx,dell", "dell,alma9"),
    ("gpu-slurm-centos78,R730",     "raptor,alma9"),
    ("groups=gpudata",              "groups=raptor"),
    ("groups=data",                 "groups=dell"),
    ("groups=vizdata",              "groups=viz"),
    ("groups=gpumgt",               "groups=mgt"),
    ("groups=vizmgmt",              "groups=vizmgt"),
    ("gpu-slurm-centos78",          "raptor"),
    ("compute-centos78-nomlx",      "dell"),
    (",R730",                       ",alma9"),

    # netboot method
    ("netboot=pxe", "netboot=xnba"),

    # OS attribute strings in commands
    ("os=centos7.8",                  "os=alma9.5"),
    ("os=centos7.6",                  "os=alma9.5"),
    ("osvers=centos7.6",              "osvers=alma9.7"),
    ("osvers=centos7.8",              "osvers=alma9.5"),
    ("osdistroname=centos7.6-x86_64", "osdistroname=alma9.7-x86_64"),
    ("osdistroname=centos7.8-x86_64", "osdistroname=alma9.5-x86_64"),

    # ISO / copycds
    ("CentOS-7-x86_64-DVD.iso",               "AlmaLinux-9.5-x86_64-dvd.iso"),
    ("-n centos7.8 -a x86_64",                "-n alma9.5 -a x86_64"),
    ("under /install/centos7.8:",             "under /install/alma9.5:"),
    ("centos7.8-x86_64-install-compute is the", "alma9.5-x86_64-install-compute is the"),
    ("centos7.8-x86_64-netboot-compute is the", "alma9.5-x86_64-netboot-compute is the"),
    ("genimage centos7.8-x86_64-netboot-panther-pec6220",
     "genimage alma9.5-x86_64-netboot-compute"),
    ("packimage centos7.8-x86_64-netboot-panther-pec6220",
     "packimage alma9.5-x86_64-netboot-compute"),

    # Management node host names
    ("ms1-dump", "hpcxcat-dump"),
    ("@ms1 ",    "@hpcxcat "),
    (" ms1 ",    " hpcxcat "),
    ("ms2",      "hpcxcat"),

    # Cluster / domain
    ("HPC Panther Cluster", "HPC Roary Cluster"),
    ("Panther Cluster",     "Roary Cluster"),
    ("panther.net",         "roary.net"),
    ("Panther",             "Roary"),

    # Management network  (DRAC 254.x and Data 233.x are unchanged)
    ("192.168.248.", "192.168.246."),

    # Example node names  (rXXX before nXXX/gXXX to avoid partial collision)
    ("rg007",      "r-a100-01"),
    ("g007",       "gpu-a100-01"),
    ("gpunode07",  "gpu-a100-01"),
    ("hpcnode125", "n126"),
    ("viznode01",  "v003"),
    ("n001-n012",  "n095-n102"),
    ("rv001",      "rv003"),
    ("v001",       "v003"),
    ("r125",       "r126"),
    ("node07 ",    "gpu-a100-01 "),

    # NIC / MAC address examples
    ("installnic=em3",                              "installnic=d4:04:e6:06:04:ac"),
    ("mac=F8:BC:12:20:16:00",                       "mac=d4:04:e6:06:04:ac"),
    ('addkcmdline="ifname=em3:F8:BC:12:20:16:00"', 'addkcmdline="ifname=eth0:d4:04:e6:06:04:ac"'),
    ("f8:bc:12:20:16:00",                           "d4:04:e6:06:04:ac"),

    # Broad OS name strings  (keep last)
    ("CentOS7.8",  "AlmaLinux 9.5"),
    ("CentOS7.6",  "AlmaLinux 9.5"),
    ("CentOS 7.8", "AlmaLinux 9.5"),
    ("CentOS 7.6", "AlmaLinux 9.5"),
    ("CentOS",     "AlmaLinux"),
    ("centos7.8",  "alma9.5"),
    ("centos7.6",  "alma9.5"),
]


# ── Slide 19: authoritative lsdef output from xcat_osimage.txt ───────────────
SLIDE19_LINES = [
    "Definition of osimage currently used for AlmaLinux 9.7 (standard compute):",
    "[root@hpcxcat ~]# lsdef -t osimage -o alma9.7-x86_64-netboot-compute",
    "Object name: alma9.7-x86_64-netboot-compute",
    "    exlist=/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.exlist",
    "    imagetype=linux",
    "    osarch=x86_64",
    "    osdistroname=alma9.7-x86_64",
    "    osname=Linux",
    "    osvers=alma9.7",
    "    otherpkgdir=/install/post/otherpkgs/alma9.7/x86_64",
    "    pkgdir=/install/alma9.7/x86_64",
    "    pkglist=/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.pkglist",
    "    postinstall=/opt/xcat/share/xcat/netboot/alma/compute.alma9.x86_64.postinstall",
    "    profile=compute",
    "    provmethod=netboot",
    "    rootimgdir=/install/netboot/alma9.7/x86_64/compute",
]


def fix_slide19(prs):
    slide = prs.slides[18]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text
        if "Object name:" in t and ("lsdef" in t or "centos7.6" in t or "alma9" in t):
            set_tf_text(shape.text_frame, SLIDE19_LINES)
            print("  [slide 19] osimage definition replaced with alma9.7 content")
            return
    print("  [slide 19] WARNING: target shape not found")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    print(f"Loading {INPUT_FILE!r} ...")
    prs = Presentation(INPUT_FILE)
    print(f"  {len(prs.slides)} slides loaded")

    print("Applying global text replacements ...")
    apply_replacements(prs, REPLACEMENTS)

    print("Rebuilding slide 19 lsdef output ...")
    fix_slide19(prs)

    print(f"Saving {OUTPUT_FILE!r} ...")
    prs.save(OUTPUT_FILE)
    print(f"Done -> {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
